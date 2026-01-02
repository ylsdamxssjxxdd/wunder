from __future__ import annotations

import os
import platform
import re
import subprocess
import sys
from functools import lru_cache
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import List, Optional, Tuple

from app.core.i18n import t

_DEFAULT_SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".html",
    ".htm",
    ".py",
    ".c",
    ".cpp",
    ".cc",
    ".h",
    ".hpp",
    ".json",
    ".js",
    ".ts",
    ".css",
    ".ini",
    ".cfg",
    ".log",
    ".doc",
    ".docx",
    ".odt",
    ".pptx",
    ".odp",
    ".xlsx",
    ".ods",
    ".wps",
    ".et",
    ".dps",
}
_DOC2MD_REF_ROOT = Path(
    os.getenv("DOC2MD_REFERENCE_ROOT", r"C:\Users\32138\Desktop\eva\thirdparty\doc2md")
)
_DOC2MD_README = _DOC2MD_REF_ROOT / "README.md"
_DOC2MD_EXT_PATTERN = re.compile(r"\*\.([a-z0-9]+)", re.IGNORECASE)
_FILENAME_SAFE_PATTERN = re.compile(r"[\\/:*?\"<>|]+")


@dataclass
class ConversionResult:
    """转换结果信息，便于调用方记录转换方式与警告。"""

    converter: str
    warnings: List[str]


class _HtmlStripper(HTMLParser):
    """轻量 HTML 文本提取器，用于兜底场景。"""

    def __init__(self) -> None:
        super().__init__()
        self._chunks: List[str] = []

    def handle_data(self, data: str) -> None:
        if data:
            self._chunks.append(data)

    def get_text(self) -> str:
        return "".join(self._chunks)


@lru_cache(maxsize=1)
def get_supported_extensions() -> List[str]:
    """读取 doc2md 说明中的扩展名列表，失败时回退内置列表。"""
    if _DOC2MD_README.exists():
        try:
            content = _DOC2MD_README.read_text(encoding="utf-8", errors="ignore")
            exts = {f".{item.lower()}" for item in _DOC2MD_EXT_PATTERN.findall(content)}
            if exts:
                return sorted(exts)
        except OSError:
            pass
    return sorted(_DEFAULT_SUPPORTED_EXTENSIONS)


def sanitize_filename_stem(name: str) -> str:
    """清理文件名，移除危险字符与路径分隔符。"""
    cleaned = _FILENAME_SAFE_PATTERN.sub("_", str(name or "").strip())
    cleaned = cleaned.strip(". ")
    cleaned = cleaned.replace("..", "_")
    return cleaned


def resolve_unique_markdown_path(root: Path, stem: str) -> Path:
    """生成不冲突的 Markdown 输出路径。"""
    safe_stem = sanitize_filename_stem(stem) or "document"
    candidate = root / f"{safe_stem}.md"
    if not candidate.exists():
        return candidate
    for idx in range(1, 1000):
        candidate = root / f"{safe_stem}-{idx}.md"
        if not candidate.exists():
            return candidate
    raise RuntimeError(t("error.converter_unique_markdown_name"))


def resolve_doc2md_binary() -> Optional[Path]:
    """根据系统平台选择可用的 doc2md 可执行文件。"""
    repo_root = Path(__file__).resolve().parents[2]
    bin_root = repo_root / "scripts" / "doc2md"
    system = sys.platform.lower()
    machine = platform.machine().lower()
    candidates: List[str] = []
    if system.startswith("win"):
        candidates.append("doc2md-win-x86_64.exe")
    elif system.startswith("linux"):
        if machine in {"aarch64", "arm64"}:
            candidates.append("doc2md-linux-arm64")
        candidates.append("doc2md-linux-x86_64")
    for name in candidates:
        path = bin_root / name
        if path.exists():
            if os.name == "posix":
                try:
                    mode = path.stat().st_mode
                    if not (mode & 0o111):
                        path.chmod(mode | 0o111)
                except OSError:
                    pass
            return path
    return None


def run_doc2md(binary: Path, input_path: Path, output_path: Path) -> Tuple[bool, str]:
    """执行 doc2md CLI，将输入文件转换为 Markdown。"""
    try:
        result = subprocess.run(
            [str(binary), "-o", str(output_path), str(input_path)],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        return False, str(exc)
    if result.returncode != 0:
        return False, (
            result.stderr or result.stdout or t("error.converter_doc2md_failed")
        ).strip()
    if not output_path.exists():
        return False, t("error.converter_doc2md_no_output")
    return True, ""


def _read_text(path: Path) -> str:
    """尝试多种编码读取文本，尽量避免失败。"""
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return path.read_text(encoding=encoding, errors="ignore")
        except OSError:
            continue
    raise RuntimeError(t("error.converter_read_text_failed"))


def _wrap_code_block(text: str, language: str) -> str:
    """将文本包裹为 Markdown 代码块。"""
    body = text.rstrip()
    return f"```{language}\n{body}\n```"


def _convert_html(text: str) -> str:
    """HTML → Markdown 的兜底转换。"""
    try:
        from markdownify import markdownify  # type: ignore

        return markdownify(text)
    except Exception:
        pass
    try:
        import html2text  # type: ignore

        handler = html2text.HTML2Text()
        handler.ignore_links = False
        return handler.handle(text)
    except Exception:
        pass
    stripper = _HtmlStripper()
    stripper.feed(text)
    return stripper.get_text()


def _convert_docx(path: Path) -> str:
    """尝试使用 python-docx 提取 docx 内容。"""
    try:
        from docx import Document  # type: ignore
    except Exception as exc:
        raise RuntimeError(
            t("error.converter_python_docx_unavailable", detail=str(exc))
        ) from exc
    doc = Document(str(path))
    lines: List[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    for table in doc.tables:
        rows: List[List[str]] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if not rows:
            continue
        header = rows[0]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines).strip()


def _convert_pptx(path: Path) -> str:
    """尝试使用 python-pptx 提取 pptx 内容。"""
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        raise RuntimeError(
            t("error.converter_python_pptx_unavailable", detail=str(exc))
        ) from exc
    presentation = Presentation(str(path))
    lines: List[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"## Slide {index}")
        texts: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = text.strip()
            if text:
                texts.append(text)
        if texts:
            for item in texts:
                lines.append(f"- {item}")
    return "\n".join(lines).strip()


def _convert_xlsx(path: Path) -> str:
    """尝试使用 openpyxl 将 xlsx 转为 Markdown 表格。"""
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:
        raise RuntimeError(
            t("error.converter_openpyxl_unavailable", detail=str(exc))
        ) from exc
    workbook = load_workbook(filename=str(path), data_only=True)
    lines: List[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"## Sheet {sheet.title}")
        rows: List[List[str]] = []
        for row in sheet.iter_rows(values_only=True):
            if row is None:
                continue
            values = ["" if cell is None else str(cell) for cell in row]
            if any(value.strip() for value in values):
                rows.append(values)
        if not rows:
            continue
        header = rows[0]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines).strip()


def _try_pypandoc(path: Path) -> Optional[str]:
    """优先尝试 pypandoc 处理多格式转换。"""
    try:
        import pypandoc  # type: ignore

        return pypandoc.convert_file(str(path), "md")
    except Exception:
        return None


def convert_with_python(path: Path, extension: str) -> Tuple[str, str]:
    """使用 Python 生态做格式转换的兜底方案。"""
    ext = extension.lower()
    markdown = _try_pypandoc(path)
    if markdown:
        return markdown, "pypandoc"

    if ext in {".md", ".markdown"}:
        return _read_text(path), "text"

    if ext in {".txt", ".log"}:
        return _read_text(path), "text"

    if ext in {".html", ".htm"}:
        return _convert_html(_read_text(path)), "html"

    code_map = {
        ".py": "python",
        ".c": "c",
        ".cpp": "cpp",
        ".cc": "cpp",
        ".h": "c",
        ".hpp": "cpp",
        ".json": "json",
        ".js": "javascript",
        ".ts": "typescript",
        ".css": "css",
        ".ini": "",
        ".cfg": "",
    }
    if ext in code_map:
        language = code_map.get(ext, "")
        return _wrap_code_block(_read_text(path), language), "code"

    if ext == ".docx":
        return _convert_docx(path), "python-docx"

    if ext == ".pptx":
        return _convert_pptx(path), "python-pptx"

    if ext == ".xlsx":
        return _convert_xlsx(path), "openpyxl"

    raise RuntimeError(t("error.converter_python_converter_not_found", ext=ext))


def convert_to_markdown(
    input_path: Path, output_path: Path, extension: str
) -> ConversionResult:
    """尝试 doc2md 与 Python 兜底转换，成功后写入 Markdown 文件。"""
    warnings: List[str] = []
    binary = resolve_doc2md_binary()
    if binary:
        ok, detail = run_doc2md(binary, input_path, output_path)
        if ok:
            return ConversionResult(converter="doc2md", warnings=warnings)
        warnings.append(detail or t("error.converter_doc2md_convert_failed"))
    markdown, converter = convert_with_python(input_path, extension)
    if not markdown.strip():
        raise RuntimeError(t("error.converter_empty_result"))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(markdown, encoding="utf-8")
    return ConversionResult(converter=converter, warnings=warnings)
