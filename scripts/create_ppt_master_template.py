from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path("config/ppt_templates/black_times_default")
PPTX_PATH = ROOT / "template.pptx"
CONFIG_PATH = ROOT / "template.json"
EAST_ASIAN_FONT = "SimHei"
LATIN_FONT = "Times New Roman"


def main() -> None:
    ROOT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _write_config()
    _write_seed_preview(prs)
    prs.save(str(PPTX_PATH))


def _write_config() -> None:
    data = {
        "id": "black_times_default",
        "name": "Black Times Default",
        "description": "Master-template PPTX pack using SimHei for Chinese and Times New Roman for Latin text.",
        "file": "template.pptx",
        "aliases": ["master_default", "simhei_times", "black_times"],
        "fonts": {
            "east_asian": EAST_ASIAN_FONT,
            "latin": LATIN_FONT,
        },
        "layouts": {
            "cover": {"index": 0, "layout_name": "Title Slide"},
            "content": {"index": 1, "layout_name": "Title and Content"},
            "toc": {"index": 1, "layout_name": "Title and Content"},
            "section": {"index": 2, "layout_name": "Section Header"},
            "content_image": {"index": 8, "layout_name": "Picture with Caption"},
            "timeline": {"index": 1, "layout_name": "Title and Content"},
            "comparison": {"index": 4, "layout_name": "Comparison"},
            "data": {"index": 1, "layout_name": "Title and Content"},
            "closing": {"index": 2, "layout_name": "Section Header"},
            "default": {"index": 1, "layout_name": "Title and Content"},
        },
    }
    CONFIG_PATH.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _write_seed_preview(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _shape(slide, 0, 0, 13.333, 7.5, "F4F6F8")
    _shape(slide, 9.2, -1.25, 5.1, 5.1, "00C9B6", shape_type=MSO_SHAPE.OVAL, transparency=12)
    _shape(slide, -1.2, 5.2, 3.4, 3.4, "10B981", shape_type=MSO_SHAPE.OVAL, transparency=16)
    _shape(slide, 0.92, 2.28, 0.11, 1.65, "00C9B6")
    _textbox(slide, "母版模板预览", 1.22, 2.08, 7.0, 0.8, 40, bold=True)
    _textbox(slide, "中文黑体 / English Times New Roman", 1.24, 3.22, 7.2, 0.42, 18)
    _textbox(slide, "第一页使用封面版式，最后一页使用结尾版式，中间页按 type/layout 选择。", 1.24, 5.95, 10.8, 0.36, 13)


def _shape(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    *,
    shape_type=MSO_SHAPE.RECTANGLE,
    transparency: int = 0,
) -> None:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = _rgb(fill)


def _textbox(slide, text: str, x: float, y: float, w: float, h: float, size: float, *, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = LATIN_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb("1F2329")
    _set_typeface(run.font._element, "a:latin", LATIN_FONT)
    _set_typeface(run.font._element, "a:ea", EAST_ASIAN_FONT)
    _set_typeface(run.font._element, "a:cs", LATIN_FONT)


def _set_typeface(r_pr, tag: str, typeface: str) -> None:
    element = r_pr.find(qn(tag))
    if element is None:
        element = OxmlElement(tag)
        r_pr.append(element)
    element.set("typeface", typeface)


def _rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


if __name__ == "__main__":
    main()
