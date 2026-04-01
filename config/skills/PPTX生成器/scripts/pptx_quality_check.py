#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
import zipfile
import re
from pathlib import Path
from typing import Any

from pptx import Presentation

EMU_PER_INCH = 914400


def emu_to_inch(value: int) -> float:
    return float(value) / EMU_PER_INCH


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Check PPTX visual-density and layout quality gates for wunder skill workflow."
    )
    parser.add_argument("pptx_path", help="Path to target pptx file.")
    parser.add_argument("--min-slides", type=int, default=8)
    parser.add_argument("--min-avg-shapes", type=float, default=10.0)
    parser.add_argument("--min-avg-text-boxes", type=float, default=6.0)
    parser.add_argument(
        "--check-page-badge",
        action="store_true",
        help="Validate non-cover slides contain badge-like elements near bottom-right.",
    )
    return parser.parse_args()


def shape_has_visible_text(shape: Any) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = shape.text_frame.text if shape.text_frame else ""
    return bool(text and text.strip())


def is_badge_zone(shape: Any) -> bool:
    # Badge is expected near (x=9.3, y=5.1, w/h around 0.4) on 16:9 slide.
    x = emu_to_inch(shape.left)
    y = emu_to_inch(shape.top)
    w = emu_to_inch(shape.width)
    h = emu_to_inch(shape.height)
    return (
        x >= 8.7
        and y >= 4.75
        and w <= 1.0
        and h <= 1.0
        and (x + w) <= 10.1
        and (y + h) <= 5.8
    )


def main() -> int:
    args = parse_args()
    target = Path(args.pptx_path).expanduser().resolve()

    if not target.exists():
        print(
            json.dumps(
                {"ok": False, "error": f"file_not_found: {target}"},
                ensure_ascii=False,
                indent=2,
            )
        )
        return 2

    prs = Presentation(str(target))
    slide_reports: list[dict[str, Any]] = []
    total_shapes = 0
    total_text_boxes = 0
    failures: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        shape_count = len(slide.shapes)
        text_count = 0
        badge_zone_shapes = 0
        badge_zone_digit_text = 0

        for shape in slide.shapes:
            if shape_has_visible_text(shape):
                text_count += 1
            if is_badge_zone(shape):
                badge_zone_shapes += 1
                if shape_has_visible_text(shape) and any(ch.isdigit() for ch in shape.text):
                    badge_zone_digit_text += 1

        total_shapes += shape_count
        total_text_boxes += text_count

        slide_report = {
            "slide": idx,
            "shape_count": shape_count,
            "text_box_count": text_count,
            "badge_zone_shapes": badge_zone_shapes,
            "badge_zone_digit_text": badge_zone_digit_text,
        }
        slide_reports.append(slide_report)

        if args.check_page_badge and idx > 1:
            # Non-cover slides should show both a badge shape and a number-like text.
            if badge_zone_shapes == 0 or badge_zone_digit_text == 0:
                failures.append(f"slide_{idx}: missing_page_badge")

    slide_count = len(prs.slides)
    avg_shapes = round(total_shapes / slide_count, 2) if slide_count else 0.0
    avg_text_boxes = round(total_text_boxes / slide_count, 2) if slide_count else 0.0

    if slide_count < args.min_slides:
        failures.append(f"slides<{args.min_slides} (actual={slide_count})")
    if avg_shapes < args.min_avg_shapes:
        failures.append(f"avg_shapes<{args.min_avg_shapes} (actual={avg_shapes})")
    if avg_text_boxes < args.min_avg_text_boxes:
        failures.append(f"avg_text_boxes<{args.min_avg_text_boxes} (actual={avg_text_boxes})")

    # Guard against a known PptxGenJS misuse: addShape("oval") -> invalid DrawingML prst token.
    invalid_prst_tokens: list[dict[str, str]] = []
    with zipfile.ZipFile(target, "r") as archive:
        for name in archive.namelist():
            if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                continue
            xml = archive.read(name).decode("utf-8", errors="ignore")
            for match in re.finditer(r'<a:prstGeom prst="([^"]+)"', xml):
                token = match.group(1)
                if token == "oval":
                    invalid_prst_tokens.append({"slide_xml": name, "token": token})
    if invalid_prst_tokens:
        failures.append(f"invalid_prst_token_detected(count={len(invalid_prst_tokens)})")

    result = {
        "ok": len(failures) == 0,
        "file": str(target),
        "summary": {
            "slides": slide_count,
            "size_bytes": target.stat().st_size,
            "total_shapes": total_shapes,
            "avg_shapes_per_slide": avg_shapes,
            "total_text_boxes": total_text_boxes,
            "avg_text_boxes_per_slide": avg_text_boxes,
            "thresholds": {
                "min_slides": args.min_slides,
                "min_avg_shapes": args.min_avg_shapes,
                "min_avg_text_boxes": args.min_avg_text_boxes,
                "check_page_badge": args.check_page_badge,
            },
        },
        "failures": failures,
        "invalid_prst_tokens": invalid_prst_tokens,
        "slides": slide_reports,
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["ok"] else 1


if __name__ == "__main__":
    sys.exit(main())
