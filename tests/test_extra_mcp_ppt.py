from __future__ import annotations

import importlib.util
import sys
import types
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def _install_mcp_stub_if_missing() -> None:
    if importlib.util.find_spec("mcp") is not None:
        return
    mcp = types.ModuleType("mcp")
    server = types.ModuleType("mcp.server")
    fastmcp = types.ModuleType("mcp.server.fastmcp")

    class FastMCP:
        def tool(self, *args, **kwargs):
            def deco(fn):
                return fn

            return deco

    fastmcp.FastMCP = FastMCP
    sys.modules["mcp"] = mcp
    sys.modules["mcp.server"] = server
    sys.modules["mcp.server.fastmcp"] = fastmcp


_install_mcp_stub_if_missing()

from extra_mcp.tools.ppt.tools import (
    _delete_sync,
    _read_sync,
    _refine_sync,
    _template_read_sync,
    _write_sync,
)


def test_ppt_tool_flow_generates_refines_reads_and_deletes(tmp_path, monkeypatch):
    workspace_root = tmp_path / "workspaces"
    workspace_root.mkdir()
    monkeypatch.setenv("EXTRA_MCP_PPT_ROOT", str(tmp_path / "ppt"))
    monkeypatch.setenv("WUNDER_WORKSPACE_ROOT", str(workspace_root))

    content = """
    <slides>
      <slide>
        <type>cover</type>
        <title>独立 PPT MCP</title>
        <subtitle>工具化生成流程</subtitle>
        <prompt>封面，商务风格，突出可编辑和可精修</prompt>
      </slide>
      <slide>
        <type>toc</type>
        <title>目录</title>
        <bullet>生成</bullet>
        <bullet>读取</bullet>
        <bullet>精修</bullet>
        <prompt>目录页，三项内容</prompt>
      </slide>
      <slide>
        <type>timeline</type>
        <title>流程</title>
        <bullet>规划页面</bullet>
        <bullet>渲染 PPTX</bullet>
        <bullet>读取 slide_id</bullet>
        <bullet>精修页面</bullet>
        <prompt>时间线流程页</prompt>
      </slide>
    </slides>
    """
    created = _write_sync(
        presentation_id="",
        presentation_name="独立 PPT MCP",
        insert_before="",
        content=content,
        lang="xml",
        template_id="research_blue",
        output_path="/workspaces/u1/exports/demo.pptx",
        overwrite=True,
    )
    assert created["ok"] is True
    assert created["slide_count"] == 3
    assert created["template_id"] == "research_blue"
    assert created["workspace_relative_path"] == "exports/demo.pptx"
    assert Path(created["output_path"]).exists()

    presentation_id = created["presentation_id"]
    refine = """
    <slides>
      <slide>
        <slide_id>slide_003</slide_id>
        <type>comparison</type>
        <title>流程优化对比</title>
        <item><title>原流程</title><body>模型直接写代码；布局不稳定</body></item>
        <item><title>新流程</title><body>模型写页面意图；工具稳定渲染</body></item>
        <prompt>改成左右对比页</prompt>
      </slide>
    </slides>
    """
    updated = _refine_sync(
        presentation_id=presentation_id,
        content=refine,
        lang="xml",
        template_id="executive_green",
        output_path="",
        overwrite=False,
    )
    assert updated["ok"] is True
    assert updated["changed_slide_ids"] == ["slide_003"]
    assert updated["template_id"] == "executive_green"

    read = _read_sync(
        presentation_id=presentation_id,
        path="",
        slide_ids=["slide_003"],
        max_slides=30,
    )
    assert read["ok"] is True
    assert read["template_id"] == "executive_green"
    assert read["slides"][0]["title"] == "流程优化对比"

    templates = _template_read_sync(template_id="", path="", max_slides=30)
    assert templates["ok"] is True
    assert templates["type"] == "builtin_template_list"
    assert {item["template_id"] for item in templates["templates"]} >= {
        "amber_clear",
        "executive_green",
        "research_blue",
        "finance_ink",
        "creative_coral",
        "minimal_gray",
    }

    template = _template_read_sync(template_id="finance", path="", max_slides=30)
    assert template["ok"] is True
    assert template["template"]["template_id"] == "finance_ink"

    deleted = _delete_sync(
        presentation_id=presentation_id,
        slide_ids=["slide_002"],
        output_path="",
        overwrite=False,
    )
    assert deleted["ok"] is True
    assert deleted["slide_count"] == 2


def test_ppt_doubao_radar_template_supports_images(tmp_path, monkeypatch):
    workspace_root = tmp_path / "workspaces"
    workspace_root.mkdir()
    monkeypatch.setenv("EXTRA_MCP_PPT_ROOT", str(tmp_path / "ppt"))
    monkeypatch.setenv("WUNDER_WORKSPACE_ROOT", str(workspace_root))

    image_path = tmp_path / "sample.png"
    Image.new("RGB", (640, 360), "#0f766e").save(image_path)
    content = f"""
    <slides>
      <slide>
        <type>content</type>
        <title>Image Support</title>
        <body>Images can be placed into the Doubao-like technical layout.</body>
        <bullet>Workspace/local image path</bullet>
        <image src="{image_path.as_posix()}" />
        <prompt>content page with right-side image</prompt>
      </slide>
    </slides>
    """
    created = _write_sync(
        presentation_id="",
        presentation_name="doubao-radar-image",
        insert_before="",
        content=content,
        lang="xml",
        template_id="doubao_radar",
        output_path="/workspaces/u1/exports/doubao-radar-image.pptx",
        overwrite=True,
    )
    assert created["ok"] is True
    assert created["slide_count"] == 1
    assert created["template_id"] == "doubao_radar"
    assert Path(created["output_path"]).exists()

    prs = Presentation(created["output_path"])
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[0].shapes)


def test_ppt_master_template_uses_layouts_and_fonts(tmp_path, monkeypatch):
    workspace_root = tmp_path / "workspaces"
    workspace_root.mkdir()
    monkeypatch.setenv("EXTRA_MCP_PPT_ROOT", str(tmp_path / "ppt"))
    monkeypatch.setenv("WUNDER_WORKSPACE_ROOT", str(workspace_root))

    templates = _template_read_sync(template_id="", path="", max_slides=30)
    master_ids = {item["template_id"] for item in templates["templates"] if item.get("type") == "master_template"}
    assert "black_times_default" in master_ids

    content = """
    <slides>
      <slide>
        <type>cover</type>
        <title>中文黑体与 Times New Roman</title>
        <subtitle>Master Template Demo</subtitle>
        <prompt>封面</prompt>
      </slide>
      <slide>
        <type>toc</type>
        <title>目录</title>
        <bullet>母版复用</bullet>
        <bullet>版式选择</bullet>
        <bullet>字体规则</bullet>
        <prompt>目录页</prompt>
      </slide>
      <slide>
        <type>comparison</type>
        <title>方案对比</title>
        <item><title>代码模板</title><body>速度快；复刻灵活</body></item>
        <item><title>母版模板</title><body>复用强；品牌一致</body></item>
        <prompt>对比页</prompt>
      </slide>
      <slide>
        <type>closing</type>
        <title>结束页</title>
        <body>最后一页强制使用结尾版式。</body>
        <prompt>结尾</prompt>
      </slide>
    </slides>
    """
    created = _write_sync(
        presentation_id="",
        presentation_name="master-demo",
        insert_before="",
        content=content,
        lang="xml",
        template_id="black_times_default",
        output_path="/workspaces/u1/exports/master-demo.pptx",
        overwrite=True,
    )
    assert created["ok"] is True
    assert created["template_id"] == "black_times_default"
    assert created["slide_count"] == 4
    assert Path(created["output_path"]).exists()

    prs = Presentation(created["output_path"])
    assert len(prs.slides) == 4
    first_text = "\n".join(
        shape.text for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False)
    )
    last_text = "\n".join(
        shape.text for shape in prs.slides[-1].shapes if getattr(shape, "has_text_frame", False)
    )
    assert "中文黑体与 Times New Roman" in first_text
    assert "结束页" in last_text

    run = _first_text_run(prs)
    assert run is not None
    fonts = run.font._element
    assert fonts.find(qn("a:latin")).get("typeface") == "Times New Roman"
    assert fonts.find(qn("a:ea")).get("typeface") == "SimHei"


def _first_text_run(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        return run
    return None
