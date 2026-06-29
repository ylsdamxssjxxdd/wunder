from __future__ import annotations

import math
import re
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .model import PresentationManifest, SlideSpec
from .templates import COLOR_KEYS, style_for_template, theme_for_template

SLIDE_W = 13.333
SLIDE_H = 7.5
SAFE_X = 0.58
SAFE_Y = 0.42
TITLE_FONT = "Microsoft YaHei"
BODY_FONT = "Microsoft YaHei"


def render_manifest(manifest: PresentationManifest, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    manifest_template_id = str(manifest.theme.get("template_id") or "")

    for index, slide in enumerate(manifest.slides, start=1):
        template_id = slide.template_id or manifest_template_id
        theme = _theme(manifest.theme, template_id=template_id)
        _render_slide(prs, slide, index, len(manifest.slides), theme, style_for_template(template_id))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def summarize_pptx(path: Path, max_slides: int = 30) -> dict[str, Any]:
    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        if index > max_slides:
            break
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        slides.append(
            {
                "slide_id": f"slide_{index:03d}",
                "index": index,
                "text": "\n".join(texts)[:2000],
                "shape_count": len(slide.shapes),
            }
        )
    return {
        "ok": True,
        "path": str(path),
        "slide_count": len(prs.slides),
        "width_in": round(prs.slide_width / 914400, 3),
        "height_in": round(prs.slide_height / 914400, 3),
        "slides": slides,
        "truncated": len(prs.slides) > max_slides,
    }


def _render_slide(
    prs: Presentation,
    spec: SlideSpec,
    index: int,
    total: int,
    theme: dict[str, str],
    style: str = "default",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if style == "doubao_radar":
        _render_doubao_slide(slide, spec, index, total, theme)
        return
    _background(slide, theme)
    slide_type = spec.slide_type
    if slide_type == "cover":
        _render_cover(slide, spec, theme)
    elif slide_type == "toc":
        _render_toc(slide, spec, index, theme)
    elif slide_type == "section":
        _render_section(slide, spec, index, theme)
    elif slide_type == "comparison":
        _render_comparison(slide, spec, theme)
    elif slide_type == "timeline":
        _render_timeline(slide, spec, theme)
    elif slide_type == "data":
        _render_data(slide, spec, theme)
    elif slide_type == "closing":
        _render_closing(slide, spec, theme)
    else:
        _render_content(slide, spec, theme)
    if index > 1:
        _page_badge(slide, index, total, theme)


def _render_doubao_slide(slide: Any, spec: SlideSpec, index: int, total: int, theme: dict[str, str]) -> None:
    _doubao_background(slide, theme)
    slide_type = spec.slide_type
    if slide_type == "cover":
        _doubao_cover(slide, spec, theme)
    elif slide_type == "toc":
        _doubao_toc(slide, spec, theme)
    elif slide_type == "timeline":
        _doubao_timeline(slide, spec, theme)
    elif slide_type == "comparison":
        _doubao_comparison(slide, spec, theme)
    elif slide_type == "data":
        _doubao_data(slide, spec, theme)
    elif slide_type == "section":
        _doubao_section(slide, spec, index, theme)
    elif slide_type == "closing":
        _doubao_closing(slide, spec, theme)
    else:
        _doubao_content(slide, spec, theme)
    if index > 1:
        _doubao_footer(slide, index, total, theme)


def _doubao_background(slide: Any, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["bg"], line=theme["bg"])
    _shape(slide, MSO_SHAPE.OVAL, 9.72, -1.39, 5.56, 5.56, theme["accent"], transparency=8)
    _shape(slide, MSO_SHAPE.OVAL, -1.39, 4.86, 4.17, 4.17, theme["accent2"], transparency=10)


def _doubao_cover(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.OVAL, 8.75, 2.15, 3.61, 3.61, theme["accent"], transparency=18)
    _shape(slide, MSO_SHAPE.OVAL, 9.38, 2.78, 2.36, 2.36, theme["bg"], line=theme["bg"], transparency=2)
    _shape(slide, MSO_SHAPE.OVAL, 10.44, 3.85, 0.22, 0.22, theme["accent"], transparency=4)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.83, 2.43, 0.11, 1.81, theme["accent"], line=theme["accent"])
    _text(slide, spec.title or "Presentation", 1.18, 2.24, 7.35, 1.65, 40, theme["primary"], bold=True)
    subtitle = spec.subtitle or spec.body or _compact_prompt(spec.prompt)
    if subtitle:
        _text(slide, subtitle, 1.18, 4.48, 7.2, 0.7, 20, theme["secondary"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.83, 6.32, 11.67, 0.01, "FFFFFF", line="FFFFFF")
    tags = _cover_tags(spec)
    if tags:
        _text(slide, " / ".join(tags), 0.83, 6.5, 5.7, 0.4, 12, theme["secondary"], bold=True)
    _text(slide, "SMART DETECTION SYSTEMS", 6.94, 6.5, 5.56, 0.4, 11, theme["muted"], align=PP_ALIGN.RIGHT)


def _doubao_toc(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _text(slide, spec.title or "目录 CONTENTS", 0.67, 0.67, 12.0, 0.64, 30, theme["primary"], bold=True)
    items = _content_items(spec, fallback_count=6)[:6]
    for idx, item in enumerate(items, start=1):
        col = 0 if idx <= 3 else 1
        row = (idx - 1) % 3
        x = 0.67 + col * 6.22
        y = 1.94 + row * 1.74
        _card(slide, x, y, 5.78, 1.46, theme, fill="FFFFFF")
        _text(slide, f"{idx:02d}", x + 0.22, y + 0.2, 1.0, 0.6, 21, theme["accent"], bold=True)
        title = str(item.get("title") or item.get("text") or f"Section {idx}")
        body = str(item.get("body") or item.get("description") or "")
        _text(slide, title, x + 1.42, y + 0.17, 4.1, 0.38, 15, theme["primary"], bold=True)
        if body:
            _text(slide, body, x + 1.42, y + 0.58, 4.05, 0.56, 8.8, theme["secondary"])


def _doubao_section(slide: Any, spec: SlideSpec, index: int, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["dark_panel"], line=theme["dark_panel"])
    _shape(slide, MSO_SHAPE.OVAL, 9.6, -1.2, 5.2, 5.2, theme["accent"], transparency=18)
    _text(slide, f"{index:02d}", 0.9, 1.08, 2.2, 1.0, 60, theme["accent"], bold=True)
    _text(slide, spec.title or "Section", 0.9, 2.55, 10.8, 0.8, 36, "FFFFFF", bold=True)
    body = spec.body or _compact_prompt(spec.prompt)
    if body:
        _text(slide, body, 0.92, 3.65, 9.8, 0.7, 16, "E5E7EB")


def _doubao_content(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _doubao_header(slide, spec.title or "Content", theme)
    image = _first_image(spec)
    if image and _image(slide, image, 7.72, 1.55, 5.05, 4.85):
        text_w = 6.35
    else:
        text_w = 6.35
        _doubao_visual_panel(slide, 7.72, 1.7, 5.05, 4.6, theme, spec.bullets or _fallback_bullets(spec))
    body = spec.body or _compact_prompt(spec.prompt)
    if body:
        _text(slide, body, 0.56, 1.82, text_w, 1.8, 14.5, theme["secondary"])
    bullets = spec.bullets or _fallback_bullets(spec)
    for idx, bullet in enumerate(bullets[:3]):
        y = 4.12 + idx * 0.88
        _card(slide, 0.56, y, text_w, 0.72, theme, fill="FFFFFF")
        _shape(slide, MSO_SHAPE.OVAL, 0.78, y + 0.18, 0.34, 0.34, theme["accent"])
        _text(slide, bullet, 1.26, y + 0.13, text_w - 1.0, 0.38, 11.5, theme["secondary"])


def _doubao_timeline(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _doubao_header(slide, spec.title or "Timeline", theme)
    items = _content_items(spec, fallback_count=4)[:4]
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.83, 2.22, 11.67, 0.02, theme["line"], line=theme["line"])
    for idx, item in enumerate(items):
        x = 0.83 + idx * 2.99
        _shape(slide, MSO_SHAPE.OVAL, x + 1.05, 2.11, 0.22, 0.22, "FFFFFF", line=theme["accent"])
        _card(slide, x, 2.78, 2.85, 3.61, theme, fill="FFFFFF")
        _text(slide, str(item.get("label") or item.get("date") or item.get("title") or f"Stage {idx+1}"), x + 0.28, 3.05, 2.3, 0.42, 17, theme["accent"], bold=True)
        title = str(item.get("subtitle") or item.get("body") or item.get("description") or "")
        parts = re_split_list(title)
        if parts:
            _text(slide, parts[0], x + 0.28, 3.72, 2.4, 0.36, 11.5, theme["primary"], bold=True)
            _text(slide, " ".join(parts[1:]) or title, x + 0.28, 4.34, 2.35, 1.7, 9.4, theme["secondary"])
    if spec.body:
        _shape(slide, MSO_SHAPE.RECTANGLE, 0.56, 6.6, 12.22, 0.01, "FFFFFF", line="FFFFFF")
        _text(slide, spec.body, 0.56, 6.74, 12.22, 0.34, 10.5, theme["secondary"], align=PP_ALIGN.CENTER)


def _doubao_comparison(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _doubao_header(slide, spec.title or "Comparison", theme)
    items = _content_items(spec, fallback_count=2)
    while len(items) < 2:
        items.append({"title": f"Option {len(items)+1}", "body": ""})
    _doubao_compare_column(slide, items[0], 0.56, 1.72, 5.9, 4.25, theme, "FFFFFF", theme["primary"])
    _doubao_compare_column(slide, items[1], 6.88, 1.72, 5.9, 4.25, theme, theme["accent"], "FFFFFF")
    summary = spec.body or _compact_prompt(spec.prompt)
    if summary:
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.56, 6.18, 12.22, 0.78, theme["accent"], line=theme["accent"])
        _text(slide, summary, 0.76, 6.3, 11.8, 0.42, 11, "FFFFFF")


def _doubao_compare_column(slide: Any, item: dict[str, Any], x: float, y: float, w: float, h: float, theme: dict[str, str], fill: str, text_color: str) -> None:
    _card(slide, x, y, w, h, theme, fill=fill)
    _text(slide, str(item.get("title") or "Option"), x + 0.27, y + 0.25, w - 0.54, 0.42, 17, text_color, bold=True)
    _shape(slide, MSO_SHAPE.RECTANGLE, x + 0.27, y + 0.88, w - 0.54, 0.01, theme["line"], line=theme["line"])
    parts = re_split_list(str(item.get("body") or item.get("description") or item.get("text") or ""))
    for idx, part in enumerate(parts[:4]):
        px = x + 0.27 + (idx % 2) * ((w - 0.78) / 2)
        py = y + 1.18 + (idx // 2) * 1.25
        card_fill = theme["surface_alt"] if fill == "FFFFFF" else theme["accent"]
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, (w - 0.9) / 2, 0.96, card_fill, line=theme["line"], transparency=0 if fill == "FFFFFF" else 12)
        _text(slide, part, px + 0.16, py + 0.13, (w - 1.25) / 2, 0.62, 9.2, text_color)


def _doubao_data(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _doubao_header(slide, spec.title or "Data Insights", theme)
    metrics = _metrics(spec)[:6]
    max_value = max([abs(float(item.get("value", 0))) for item in metrics] or [1])
    for idx, item in enumerate(metrics):
        col = idx % 3
        row = idx // 3
        x = 0.56 + col * 4.17
        y = 1.65 + row * 1.7
        _card(slide, x, y, 3.9, 1.35, theme, fill=theme["dark_panel"])
        _text(slide, str(item.get("label") or f"Metric {idx+1}"), x + 0.24, y + 0.18, 2.6, 0.28, 12, "FFFFFF", bold=True)
        value = float(item.get("value") or 0)
        _text(slide, _format_number(value), x + 2.9, y + 0.12, 0.72, 0.34, 16, theme["accent"], bold=True, align=PP_ALIGN.RIGHT)
        _shape(slide, MSO_SHAPE.RECTANGLE, x + 0.25, y + 0.85, 3.35, 0.12, "FFFFFF", line="FFFFFF", transparency=50)
        _shape(slide, MSO_SHAPE.RECTANGLE, x + 0.25, y + 0.85, max(0.1, 3.35 * abs(value) / max_value), 0.12, theme["accent"], line=theme["accent"])
    bullets = spec.bullets or _fallback_bullets(spec)
    _card(slide, 0.56, 5.18, 12.22, 1.55, theme, fill="FFFFFF")
    _text(slide, spec.body or "Conclusion", 0.86, 5.42, 2.2, 0.35, 16, theme["accent"], bold=True)
    _text(slide, "  ".join(bullets[:3]), 0.86, 5.95, 11.4, 0.48, 11.5, theme["secondary"])


def _doubao_closing(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["dark_panel"], line=theme["dark_panel"])
    _shape(slide, MSO_SHAPE.OVAL, 8.2, -1.0, 5.0, 5.0, theme["accent"], transparency=18)
    _text(slide, spec.title or "Thank You", 1.0, 2.4, 11.2, 0.8, 44, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    body = spec.body or _compact_prompt(spec.prompt)
    if body:
        _text(slide, body, 2.0, 3.45, 9.3, 0.5, 16, "E5E7EB", align=PP_ALIGN.CENTER)


def _doubao_header(slide: Any, title: str, theme: dict[str, str]) -> None:
    _text(slide, title, 0.56, 0.56, 12.22, 0.64, 27, theme["primary"], bold=True)


def _doubao_footer(slide: Any, index: int, total: int, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 11.6, 6.94, 0.88, 0.02, theme["accent"], line=theme["accent"])
    _text(slide, f"{index:02d}/{total:02d}", 11.72, 6.72, 0.75, 0.22, 8, theme["muted"], align=PP_ALIGN.RIGHT)


def _doubao_visual_panel(slide: Any, x: float, y: float, w: float, h: float, theme: dict[str, str], bullets: list[str]) -> None:
    _card(slide, x, y, w, h, theme, fill="FFFFFF")
    cx = x + w / 2
    cy = y + h / 2
    _shape(slide, MSO_SHAPE.OVAL, cx - 1.35, cy - 1.35, 2.7, 2.7, theme["accent"], transparency=12)
    _shape(slide, MSO_SHAPE.OVAL, cx - 0.82, cy - 0.82, 1.64, 1.64, theme["bg"], line=theme["bg"])
    _shape(slide, MSO_SHAPE.OVAL, cx - 0.12, cy - 0.12, 0.24, 0.24, theme["accent"], transparency=5)
    for idx, bullet in enumerate(bullets[:4]):
        angle = math.radians(idx * 90 + 35)
        px = cx + math.cos(angle) * 1.8
        py = cy + math.sin(angle) * 1.3
        _shape(slide, MSO_SHAPE.OVAL, px - 0.12, py - 0.12, 0.24, 0.24, theme["accent2"])


def _render_cover(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 4.2, SLIDE_H, theme["cover_panel"], line=theme["cover_panel"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 4.2, 0, 0.18, SLIDE_H, theme["accent"], line=theme["accent"])
    _shape(slide, MSO_SHAPE.OVAL, 10.0, 0.7, 2.0, 2.0, theme["accent2"], transparency=16)
    _shape(slide, MSO_SHAPE.OVAL, 11.05, 4.6, 1.15, 1.15, theme["accent"], transparency=18)
    _shape(slide, MSO_SHAPE.RECTANGLE, 9.0, 3.0, 3.6, 0.16, theme["line"], transparency=12)
    _text(slide, spec.title or "Presentation", 4.75, 1.65, 7.7, 1.25, 42, theme["primary"], bold=True)
    subtitle = spec.subtitle or spec.body or _compact_prompt(spec.prompt)
    if subtitle:
        _text(slide, subtitle, 4.78, 3.0, 6.8, 0.85, 18, theme["secondary"])
    for i, word in enumerate(_cover_tags(spec)):
        _pill(slide, word, 4.78 + i * 1.35, 4.35, 1.08, 0.36, theme["surface"], theme["accent2"], 10)
    _text(slide, "PRESENTATION", 0.55, 6.42, 3.0, 0.35, 10, "FFFFFF", bold=True)


def _render_toc(slide: Any, spec: SlideSpec, index: int, theme: dict[str, str]) -> None:
    _header(slide, spec.title or "Agenda", "Agenda", theme)
    items = _content_items(spec, fallback_count=6)
    cols = 2 if len(items) > 3 else 1
    card_w = 5.7 if cols == 2 else 11.2
    card_h = 0.82
    start_y = 1.55
    gap_y = 0.28
    for idx, item in enumerate(items[:6], start=1):
        col = (idx - 1) % cols
        row = (idx - 1) // cols
        x = SAFE_X + col * 6.05
        y = start_y + row * (card_h + gap_y)
        _card(slide, x, y, card_w, card_h, theme, fill=theme["surface"])
        _shape(slide, MSO_SHAPE.OVAL, x + 0.22, y + 0.2, 0.42, 0.42, theme["accent"])
        _text(slide, f"{idx:02d}", x + 0.22, y + 0.27, 0.42, 0.2, 9, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _text(slide, item.get("title") or item.get("text") or f"Section {idx}", x + 0.82, y + 0.16, card_w - 1.0, 0.28, 17, theme["primary"], bold=True)
        body = str(item.get("body") or item.get("description") or "")
        if body:
            _text(slide, body, x + 0.82, y + 0.49, card_w - 1.0, 0.22, 10, theme["muted"])
    _footer_line(slide, theme, index)


def _render_section(slide: Any, spec: SlideSpec, index: int, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["section_bg"], line=theme["section_bg"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.42, SLIDE_H, theme["accent"], line=theme["accent"])
    _text(slide, f"{index:02d}", 1.0, 1.15, 2.2, 1.2, 68, theme["accent"], bold=True)
    _text(slide, spec.title or "Section", 1.0, 2.65, 9.8, 0.9, 38, "FFFFFF", bold=True)
    body = spec.body or _compact_prompt(spec.prompt)
    if body:
        _text(slide, body, 1.02, 3.78, 8.6, 0.62, 17, "E5E7EB")
    _shape(slide, MSO_SHAPE.RECTANGLE, 9.5, 5.7, 2.6, 0.18, theme["accent2"], transparency=10)


def _render_content(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _header(slide, spec.title or "Content", spec.subtitle, theme)
    bullets = spec.bullets or _fallback_bullets(spec)
    left_w = 5.25
    _card(slide, SAFE_X, 1.45, left_w, 5.15, theme)
    _text(slide, spec.body or "Key Points", SAFE_X + 0.34, 1.78, left_w - 0.68, 0.46, 18, theme["primary"], bold=True)
    for idx, bullet in enumerate(bullets[:5], start=1):
        y = 2.35 + (idx - 1) * 0.76
        _shape(slide, MSO_SHAPE.OVAL, SAFE_X + 0.36, y + 0.02, 0.28, 0.28, theme["accent2"])
        _text(slide, str(idx), SAFE_X + 0.36, y + 0.065, 0.28, 0.12, 7, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _text(slide, bullet, SAFE_X + 0.78, y - 0.02, left_w - 1.04, 0.42, 13.5, theme["secondary"])
    right_x = 6.25
    _visual_grid(slide, right_x, 1.45, 6.45, 5.15, theme, bullets)


def _render_comparison(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _header(slide, spec.title or "Comparison", spec.subtitle, theme)
    items = _content_items(spec, fallback_count=2)
    if len(items) < 2:
        bullets = spec.bullets or _fallback_bullets(spec)
        mid = max(1, math.ceil(len(bullets) / 2))
        items = [
            {"title": "Option A", "body": "; ".join(bullets[:mid])},
            {"title": "Option B", "body": "; ".join(bullets[mid:])},
        ]
    cols = min(3, max(2, len(items[:3])))
    w = (11.8 - (cols - 1) * 0.34) / cols
    for idx, item in enumerate(items[:3]):
        x = SAFE_X + idx * (w + 0.34)
        accent = [theme["accent"], theme["accent2"], theme["accent3"]][idx % 3]
        _card(slide, x, 1.45, w, 5.15, theme)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, 1.45, w, 0.14, accent, line=accent)
        _text(slide, str(item.get("title") or item.get("label") or f"Option {idx + 1}"), x + 0.28, 1.82, w - 0.56, 0.4, 19, theme["primary"], bold=True)
        body = str(item.get("body") or item.get("description") or item.get("text") or "")
        parts = [part.strip() for part in re_split_list(body) if part.strip()] or [str(item.get("text") or "")]
        for line_idx, part in enumerate(parts[:5]):
            y = 2.55 + line_idx * 0.55
            _shape(slide, MSO_SHAPE.OVAL, x + 0.3, y + 0.06, 0.16, 0.16, accent)
            _text(slide, part, x + 0.56, y - 0.01, w - 0.86, 0.32, 12, theme["secondary"])


def _render_timeline(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _header(slide, spec.title or "Timeline", spec.subtitle, theme)
    items = _content_items(spec, fallback_count=5)
    count = min(5, max(3, len(items)))
    start_x = 0.86
    step_w = 2.22
    y = 3.0
    _shape(slide, MSO_SHAPE.RECTANGLE, 1.1, y + 0.24, 10.9, 0.08, theme["line"], line=theme["line"])
    for idx in range(count):
        item = items[idx] if idx < len(items) else {"title": f"Step {idx + 1}", "body": ""}
        x = start_x + idx * step_w
        accent = [theme["accent"], theme["accent2"], theme["accent3"], theme["success"], theme["danger"]][idx % 5]
        _shape(slide, MSO_SHAPE.OVAL, x, y, 0.58, 0.58, accent)
        _text(slide, str(idx + 1), x, y + 0.1, 0.58, 0.22, 12, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _text(slide, str(item.get("title") or item.get("text") or f"Step {idx + 1}"), x - 0.22, y + 0.82, 1.45, 0.42, 15, theme["primary"], bold=True, align=PP_ALIGN.CENTER)
        body = str(item.get("body") or item.get("description") or "")
        if body:
            _text(slide, body, x - 0.34, y + 1.32, 1.7, 0.58, 10.5, theme["muted"], align=PP_ALIGN.CENTER)


def _render_data(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _header(slide, spec.title or "Data Insights", spec.subtitle, theme)
    metrics = _metrics(spec)
    max_value = max([abs(float(item.get("value", 0))) for item in metrics] or [1])
    _card(slide, SAFE_X, 1.45, 7.15, 5.15, theme)
    _text(slide, "Key Metrics", 0.92, 1.78, 2.5, 0.36, 18, theme["primary"], bold=True)
    for idx, item in enumerate(metrics[:6]):
        y = 2.35 + idx * 0.55
        label = str(item.get("label") or f"Metric {idx + 1}")[:20]
        value = float(item.get("value") or 0)
        width = 4.2 * (abs(value) / max_value if max_value else 0.1)
        _text(slide, label, 0.95, y, 1.55, 0.24, 10.5, theme["secondary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, 2.65, y + 0.04, 4.25, 0.18, theme["line"], line=theme["line"])
        _shape(slide, MSO_SHAPE.RECTANGLE, 2.65, y + 0.04, max(0.08, width), 0.18, [theme["accent"], theme["accent2"], theme["accent3"]][idx % 3])
        _text(slide, _format_number(value), 7.0, y - 0.02, 0.5, 0.25, 10.5, theme["primary"], bold=True)
    _card(slide, 8.1, 1.45, 4.6, 5.15, theme, fill=theme["data_panel"])
    _text(slide, "Conclusion", 8.45, 1.86, 1.0, 0.32, 17, theme["accent"], bold=True)
    bullets = spec.bullets or _fallback_bullets(spec)
    for idx, bullet in enumerate(bullets[:4]):
        _text(slide, f"{idx + 1}. {bullet}", 8.45, 2.38 + idx * 0.68, 3.65, 0.4, 12.5, theme["secondary"])


def _render_closing(slide: Any, spec: SlideSpec, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["section_bg"], line=theme["section_bg"])
    _shape(slide, MSO_SHAPE.OVAL, 9.8, 0.8, 2.1, 2.1, theme["accent"], transparency=22)
    _shape(slide, MSO_SHAPE.OVAL, 1.2, 5.2, 1.2, 1.2, theme["accent2"], transparency=18)
    _text(slide, spec.title or "Thank You", 1.0, 1.55, 8.5, 0.9, 44, "FFFFFF", bold=True)
    bullets = spec.bullets or _fallback_bullets(spec)
    for idx, bullet in enumerate(bullets[:4]):
        _shape(slide, MSO_SHAPE.OVAL, 1.05, 3.05 + idx * 0.56, 0.18, 0.18, theme["accent"])
        _text(slide, bullet, 1.38, 2.96 + idx * 0.56, 7.8, 0.3, 14, "E5E7EB")
    if spec.body:
        _text(slide, spec.body, 1.02, 5.92, 8.8, 0.4, 13, "CBD5E1")


def _header(slide: Any, title: str, subtitle: str, theme: dict[str, str]) -> None:
    _text(slide, title, SAFE_X, 0.45, 10.5, 0.55, 28, theme["primary"], bold=True)
    if subtitle:
        _text(slide, subtitle, SAFE_X + 0.02, 1.02, 8.3, 0.28, 11.5, theme["muted"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 11.2, 0.57, 1.2, 0.12, theme["accent"], line=theme["accent"])


def _background(slide: Any, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme["bg"], line=theme["bg"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.08, theme["accent"], line=theme["accent"])


def _visual_grid(slide: Any, x: float, y: float, w: float, h: float, theme: dict[str, str], bullets: list[str]) -> None:
    colors = [theme["accent"], theme["accent2"], theme["accent3"], theme["success"]]
    _card(slide, x, y, w, h, theme, fill=theme["surface_alt"])
    _text(slide, "Structure View", x + 0.35, y + 0.3, 2.0, 0.3, 16, theme["primary"], bold=True)
    positions = [(x + 0.45, y + 1.12), (x + 3.35, y + 1.12), (x + 0.45, y + 3.1), (x + 3.35, y + 3.1)]
    for idx, (px, py) in enumerate(positions):
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, 2.35, 1.38, "FFFFFF", line=theme["line"])
        _shape(slide, MSO_SHAPE.OVAL, px + 0.22, py + 0.22, 0.35, 0.35, colors[idx])
        label = bullets[idx] if idx < len(bullets) else f"Module {idx + 1}"
        _text(slide, label, px + 0.24, py + 0.75, 1.8, 0.36, 12, theme["secondary"], bold=idx == 0)


def _footer_line(slide: Any, theme: dict[str, str], index: int) -> None:
    _text(slide, f"SECTION {index:02d}", SAFE_X, 6.76, 1.6, 0.2, 8, theme["muted"])
    _shape(slide, MSO_SHAPE.RECTANGLE, 2.0, 6.85, 9.8, 0.02, theme["line"], line=theme["line"])


def _page_badge(slide: Any, index: int, total: int, theme: dict[str, str]) -> None:
    _shape(slide, MSO_SHAPE.OVAL, 12.28, 6.78, 0.38, 0.38, theme["accent"])
    _text(slide, str(index), 12.28, 6.85, 0.38, 0.15, 8, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)


def _card(slide: Any, x: float, y: float, w: float, h: float, theme: dict[str, str], fill: str | None = None) -> None:
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill or theme["surface"], line=theme["line"])


def _pill(slide: Any, text: str, x: float, y: float, w: float, h: float, fill: str, color: str, font_size: float) -> None:
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line=color)
    _text(slide, text, x, y + 0.08, w, h - 0.1, font_size, color, bold=True, align=PP_ALIGN.CENTER)


def _image(slide: Any, path_text: str, x: float, y: float, w: float, h: float) -> Any | None:
    path = _resolve_image_path(path_text)
    if not path:
        return None
    try:
        with Image.open(path) as image:
            image_w, image_h = image.size
        shape = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
        actual_h = shape.height / 914400
        if actual_h > h:
            slide.shapes._spTree.remove(shape._element)
            shape = slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
            actual_w = shape.width / 914400
            if actual_w > w:
                ratio = actual_w / w
                shape.crop_left = max(0, min(0.45, (1 - 1 / ratio) / 2))
                shape.crop_right = shape.crop_left
                shape.width = Inches(w)
        shape.left = Inches(x)
        shape.top = Inches(y)
        if image_w and image_h:
            shape.name = f"image_{path.name}"
        return shape
    except Exception:
        return None


def _shape(
    slide: Any,
    shape_type: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    *,
    line: str | None = None,
    transparency: int = 0,
) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    shape.line.color.rgb = _rgb(line or fill)
    return shape


def _text(
    slide: Any,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: float,
    color: str,
    *,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(value or "")
    run.font.name = TITLE_FONT if bold else BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _theme(theme: dict[str, Any], *, template_id: str = "") -> dict[str, str]:
    output = theme_for_template(template_id)
    if isinstance(theme, dict):
        for key, value in theme.items():
            if key in COLOR_KEYS and isinstance(value, str):
                cleaned = value.strip().lstrip("#")
                if len(cleaned) == 6:
                    output[key] = cleaned.upper()
    return output


def _rgb(value: str) -> RGBColor:
    cleaned = value.strip().lstrip("#")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def _content_items(spec: SlideSpec, fallback_count: int) -> list[dict[str, Any]]:
    if spec.items:
        return spec.items
    if spec.sections:
        return spec.sections
    bullets = spec.bullets or _fallback_bullets(spec)
    return [{"title": bullet, "body": ""} for bullet in bullets[:fallback_count]]


def _fallback_bullets(spec: SlideSpec) -> list[str]:
    if spec.body:
        parts = re_split_list(spec.body)
        if len(parts) >= 2:
            return parts[:6]
    prompt = _compact_prompt(spec.prompt)
    if prompt:
        parts = re_split_list(prompt)
        return parts[:6] if len(parts) >= 2 else [prompt[:80]]
    return ["Define goals", "Map the path", "Highlight focus", "Form conclusion"]


def _metrics(spec: SlideSpec) -> list[dict[str, Any]]:
    if spec.metrics:
        return spec.metrics
    bullets = spec.bullets or _fallback_bullets(spec)
    return [{"label": bullet[:12], "value": 30 + idx * 12} for idx, bullet in enumerate(bullets[:5])]


def _cover_tags(spec: SlideSpec) -> list[str]:
    candidates = spec.bullets[:3] if spec.bullets else []
    if not candidates and spec.slide_type:
        candidates = ["Clear", "Professional", "Editable"]
    return [item[:6] for item in candidates[:3]]


def _compact_prompt(prompt: str) -> str:
    return " ".join(str(prompt or "").split())[:180]


def _format_number(value: float) -> str:
    if abs(value - int(value)) < 0.001:
        return str(int(value))
    return f"{value:.1f}"


def re_split_list(text: str) -> list[str]:
    return [part.strip() for part in re.split(r"[;；。、|,\n]+", str(text or "")) if part.strip()]


def _first_image(spec: SlideSpec) -> str:
    for item in spec.images:
        source = str(item.get("src") or item.get("path") or "").strip()
        if source:
            return source
    return ""


def _resolve_image_path(path_text: str) -> Path | None:
    raw = str(path_text or "").strip()
    if not raw or raw.startswith(("http://", "https://")):
        return None
    path = Path(raw).expanduser()
    candidates = [path]
    if not path.is_absolute():
        candidates.extend(
            [
                Path.cwd() / path,
                Path("/workspaces") / raw.lstrip("/\\"),
            ]
        )
    for candidate in candidates:
        try:
            resolved = candidate.resolve()
        except OSError:
            continue
        if resolved.exists() and resolved.is_file() and resolved.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff"}:
            return resolved
    return None
