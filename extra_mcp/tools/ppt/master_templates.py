from __future__ import annotations

import json
from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .model import PresentationManifest, SlideSpec
from .render import (
    BODY_FONT,
    SLIDE_H,
    SLIDE_W,
    TITLE_FONT,
    _first_image,
    _format_number,
    _image,
    _compact_prompt,
    _resolve_image_path,
    re_split_list,
)

DEFAULT_TEMPLATE_ROOT = Path("config/ppt_templates")


@dataclass(frozen=True)
class FontSpec:
    east_asian: str = "SimHei"
    latin: str = "Times New Roman"


@dataclass(frozen=True)
class LayoutSpec:
    name: str = ""
    index: int | None = None
    role: str = ""


@dataclass(frozen=True)
class MasterTemplate:
    template_id: str
    name: str
    description: str
    root: Path
    pptx_path: Path
    fill_mode: str = ""
    fonts: FontSpec = field(default_factory=FontSpec)
    layouts: dict[str, LayoutSpec] = field(default_factory=dict)
    aliases: tuple[str, ...] = ()


def template_root() -> Path:
    return Path.cwd() / DEFAULT_TEMPLATE_ROOT


def is_master_template(template_id: str | None) -> bool:
    return load_master_template(template_id) is not None


def load_master_template(template_id: str | None) -> MasterTemplate | None:
    raw = (template_id or "").strip()
    if not raw:
        return None
    normalized = _normalize_id(raw)
    for template_dir in _template_dirs():
        config_path = _template_config_path(template_dir)
        if not config_path:
            continue
        data = _read_config(config_path)
        found_id = _normalize_id(str(data.get("id") or template_dir.name))
        aliases = tuple(_normalize_id(str(item)) for item in _list(data.get("aliases")))
        if normalized not in {found_id, *aliases}:
            continue
        pptx_name = str(data.get("file") or data.get("pptx") or "template.pptx")
        pptx_path = (template_dir / pptx_name).resolve()
        if not pptx_path.exists():
            continue
        fonts = data.get("fonts") if isinstance(data.get("fonts"), dict) else {}
        return MasterTemplate(
            template_id=found_id,
            name=str(data.get("name") or found_id),
            description=str(data.get("description") or ""),
            root=template_dir,
            pptx_path=pptx_path,
            fill_mode=str(data.get("fill_mode") or data.get("mode") or "").strip().lower(),
            fonts=FontSpec(
                east_asian=str(fonts.get("east_asian") or fonts.get("ea") or "SimHei"),
                latin=str(fonts.get("latin") or "Times New Roman"),
            ),
            layouts=_parse_layouts(data.get("layouts")),
            aliases=aliases,
        )
    return None


def list_master_templates() -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    for template_dir in _template_dirs():
        config_path = _template_config_path(template_dir)
        if not config_path:
            continue
        data = _read_config(config_path)
        template_id = _normalize_id(str(data.get("id") or template_dir.name))
        template = load_master_template(template_id)
        if template is None:
            continue
        output.append(master_template_summary(template.template_id))
    return output


def master_template_summary(template_id: str) -> dict[str, Any]:
    template = load_master_template(template_id)
    if template is None:
        raise ValueError(f"master template not found: {template_id}")
    return {
        "template_id": template.template_id,
        "name": template.name,
        "description": template.description,
        "type": "master_template",
        "path": str(template.pptx_path),
        "fill_mode": template.fill_mode,
        "fonts": {
            "east_asian": template.fonts.east_asian,
            "latin": template.fonts.latin,
        },
        "aliases": list(template.aliases),
        "layouts": {
            key: {
                "name": value.name,
                "index": value.index,
                "role": value.role,
            }
            for key, value in template.layouts.items()
        },
    }


def normalize_master_template_id(template_id: str | None) -> str | None:
    template = load_master_template(template_id)
    return template.template_id if template else None


def render_master_manifest(manifest: PresentationManifest, output_path: Path, template_id: str) -> None:
    template = load_master_template(template_id)
    if template is None:
        raise ValueError(f"master template not found: {template_id}")
    prs = Presentation(str(template.pptx_path))
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    seed_slide_count = len(prs.slides)

    for index, spec in enumerate(manifest.slides, start=1):
        role = _slide_role(spec, index, len(manifest.slides))
        layout = _select_layout(prs, template, spec, role)
        slide = prs.slides.add_slide(layout)
        _fill_slide(slide, spec, index, len(manifest.slides), role, template)

    _remove_template_seed_slides(prs, seed_slide_count)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _fill_slide(slide: Any, spec: SlideSpec, index: int, total: int, role: str, template: MasterTemplate) -> None:
    if template.fill_mode == "top_title_section":
        _fill_top_title_section_slide(slide, spec, index, total, role, template)
        return

    placeholders = _placeholders(slide)
    _set_placeholder_text(placeholders, {"title", "center_title"}, spec.title, template.fonts, bold=True)
    _set_placeholder_text(placeholders, {"subtitle"}, spec.subtitle or spec.body, template.fonts)
    body_text = _body_text(spec, role)
    _set_placeholder_text(placeholders, {"body", "content", "object"}, body_text, template.fonts)
    _set_placeholder_text(placeholders, {"footer"}, f"{index:02d}/{total:02d}", template.fonts)

    if not _place_image(slide, placeholders, _first_image(spec)):
        image = _first_image(spec)
        if image:
            _image(slide, image, 7.55, 1.55, 5.0, 4.55)

    if role in {"timeline", "comparison", "data"}:
        _append_structured_text(slide, spec, role, template.fonts)


def _fill_top_title_section_slide(
    slide: Any,
    spec: SlideSpec,
    index: int,
    total: int,
    role: str,
    template: MasterTemplate,
) -> None:
    placeholders = _placeholders(slide)
    title = spec.title or template.name
    section_title = spec.subtitle or _top_section_title(spec, role, index, total)
    _set_placeholder_text(placeholders, {"title", "center_title"}, title, template.fonts, inherit_style=True)
    _set_placeholder_text(placeholders, {"subtitle"}, section_title, template.fonts, inherit_style=True)
    _set_placeholder_text(placeholders, {"footer"}, f"{index:02d}/{total:02d}", template.fonts)

    image = _first_image(spec)
    has_image = bool(image and _image(slide, image, 7.52, 2.22, 4.82, 3.62))
    content_w = 6.55 if has_image else 10.92
    content_text = _top_section_body_text(spec, role)
    if content_text:
        box = slide.shapes.add_textbox(Inches(1.16), Inches(2.24), Inches(content_w), Inches(3.65))
        _set_text_frame(box.text_frame, content_text, template.fonts, size=18 if role == "cover" else 15, bold=False)

    if role == "data" and spec.metrics:
        _append_top_section_metrics(slide, spec, template.fonts, has_image)
    elif role in {"timeline", "comparison"} and (spec.items or spec.bullets):
        _append_top_section_items(slide, spec, template.fonts, has_image)


def _top_section_title(spec: SlideSpec, role: str, index: int, total: int) -> str:
    if role == "cover":
        return spec.body or spec.prompt or "Overview"
    if role == "closing":
        return "Summary"
    if role == "toc":
        return "Contents"
    if role == "section":
        return f"Section {index:02d}"
    if total > 1:
        return f"{index:02d}"
    return "Key Point"


def _top_section_body_text(spec: SlideSpec, role: str) -> str:
    if role == "closing":
        parts = [spec.body or _compact_prompt(spec.prompt)]
        parts.extend(spec.bullets[:4])
        return "\n".join(part for part in parts if part)
    if role == "toc":
        items = _content_items_for_top_section(spec, 8)
        return "\n".join(
            f"{idx:02d}  {item.get('title') or item.get('text') or item.get('label') or ''}".rstrip()
            for idx, item in enumerate(items, start=1)
        )
    return _body_text(spec, role)


def _content_items_for_top_section(spec: SlideSpec, fallback_count: int) -> list[dict[str, Any]]:
    if spec.items:
        return spec.items
    if spec.sections:
        return spec.sections
    bullets = spec.bullets or []
    if not bullets and spec.body:
        bullets = [part.strip() for part in re_split_list(spec.body) if part.strip()]
    return [{"title": item, "body": ""} for item in bullets[:fallback_count]]


def _append_top_section_items(slide: Any, spec: SlideSpec, fonts: FontSpec, has_image: bool) -> None:
    items = _content_items_for_top_section(spec, 4)[:4]
    if not items:
        return
    x = 1.16
    y = 5.95
    width = 6.55 if has_image else 10.92
    text = "   ".join(str(item.get("title") or item.get("label") or item.get("text") or "") for item in items)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.72))
    _set_text_frame(box.text_frame, text, fonts, size=12, bold=False)


def _append_top_section_metrics(slide: Any, spec: SlideSpec, fonts: FontSpec, has_image: bool) -> None:
    metrics = spec.metrics[:4]
    if not metrics:
        return
    x = 1.16
    y = 5.86
    width = 6.55 if has_image else 10.92
    text = "   ".join(f"{item.get('label', '')}: {_metric_text(item)}" for item in metrics)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.72))
    _set_text_frame(box.text_frame, text, fonts, size=12, bold=True)


def _append_structured_text(slide: Any, spec: SlideSpec, role: str, fonts: FontSpec) -> None:
    if role == "data":
        text = "\n".join(f"{item.get('label', '')}: {item.get('value', '')}" for item in spec.metrics[:6])
    elif spec.items:
        text = "\n".join(
            f"{item.get('label') or item.get('title') or ''}  {item.get('body') or item.get('description') or item.get('text') or ''}".strip()
            for item in spec.items[:6]
        )
    else:
        text = "\n".join(spec.bullets[:6])
    if not text:
        return
    box = slide.shapes.add_textbox(Inches(0.72), Inches(5.66), Inches(11.85), Inches(0.78))
    _set_text_frame(box.text_frame, text, fonts, size=12, bold=False)


def _body_text(spec: SlideSpec, role: str) -> str:
    parts: list[str] = []
    if spec.body:
        parts.append(spec.body)
    if role == "data" and spec.metrics:
        parts.extend(f"{item.get('label', '')}: {_metric_text(item)}" for item in spec.metrics[:6])
    elif spec.items:
        parts.extend(
            f"{item.get('title') or item.get('label') or ''}: {item.get('body') or item.get('description') or item.get('text') or ''}".strip(": ")
            for item in spec.items[:6]
        )
    else:
        parts.extend(spec.bullets[:6])
    return "\n".join(part for part in parts if part)


def _metric_text(item: dict[str, Any]) -> str:
    value = item.get("value")
    try:
        return _format_number(float(value or 0))
    except (TypeError, ValueError):
        return str(value or "")


def _set_placeholder_text(
    placeholders: dict[str, list[Any]],
    roles: set[str],
    value: str,
    fonts: FontSpec,
    *,
    bold: bool = False,
    inherit_style: bool = False,
) -> bool:
    text = str(value or "").strip()
    if not text:
        return False
    for role in roles:
        for shape in placeholders.get(role, []):
            if not getattr(shape, "has_text_frame", False):
                continue
            if inherit_style:
                _set_placeholder_inherited_text(shape, text)
                return True
            size = 34 if role in {"title", "center_title"} else 15
            _set_text_frame(shape.text_frame, text, fonts, size=size, bold=bold)
            return True
    return False


def _set_placeholder_inherited_text(shape: Any, value: str) -> None:
    # Preserve the placeholder's master/layout text style by avoiding explicit run formatting.
    shape.text = value


def _set_text_frame(text_frame: Any, value: str, fonts: FontSpec, *, size: float, bold: bool) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    lines = str(value or "").splitlines() or [""]
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = line
        _apply_font(run, fonts, size=size, bold=bold)


def _apply_font(run: Any, fonts: FontSpec, *, size: float, bold: bool) -> None:
    run.font.name = fonts.latin
    run.font.size = Pt(size)
    run.font.bold = bold
    _set_typeface(run.font._element, "a:latin", fonts.latin)
    _set_typeface(run.font._element, "a:ea", fonts.east_asian)
    _set_typeface(run.font._element, "a:cs", fonts.latin)


def _set_typeface(r_pr: Any, tag: str, typeface: str) -> None:
    element = r_pr.find(qn(tag))
    if element is None:
        element = OxmlElement(tag)
        r_pr.append(element)
    element.set("typeface", typeface)


def _place_image(slide: Any, placeholders: dict[str, list[Any]], source: str) -> bool:
    path = _resolve_image_path(source)
    if path is None:
        return False
    candidates = placeholders.get("picture", []) + placeholders.get("object", [])
    for shape in candidates:
        if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(str(path), left, top, width=width, height=height)
            return True
        except Exception:
            continue
    return False


def _placeholders(slide: Any) -> dict[str, list[Any]]:
    output: dict[str, list[Any]] = {}
    for shape in slide.placeholders:
        role = _placeholder_role(shape)
        output.setdefault(role, []).append(shape)
    return output


def _placeholder_role(shape: Any) -> str:
    try:
        ph_type = shape.placeholder_format.type
    except Exception:
        return "unknown"
    if ph_type == PP_PLACEHOLDER.CENTER_TITLE:
        return "center_title"
    if ph_type == PP_PLACEHOLDER.TITLE:
        return "title"
    if ph_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if ph_type == PP_PLACEHOLDER.BODY:
        return "body"
    if ph_type == PP_PLACEHOLDER.PICTURE:
        return "picture"
    if ph_type == PP_PLACEHOLDER.OBJECT:
        return "object"
    if ph_type == PP_PLACEHOLDER.FOOTER:
        return "footer"
    if ph_type == PP_PLACEHOLDER.SLIDE_NUMBER:
        return "slide_number"
    return "unknown"


def _select_layout(prs: Presentation, template: MasterTemplate, spec: SlideSpec, role: str) -> Any:
    keys = [spec.layout, spec.template_slide_id, role, spec.slide_type, "content", "default"]
    for key in keys:
        layout_spec = template.layouts.get(str(key or "").strip().lower())
        if not layout_spec:
            continue
        layout = _find_layout(prs, layout_spec)
        if layout is not None:
            return layout
    fallback_index = 1 if len(prs.slide_layouts) > 1 else 0
    return prs.slide_layouts[fallback_index]


def _find_layout(prs: Presentation, spec: LayoutSpec) -> Any | None:
    if spec.index is not None and 0 <= spec.index < len(prs.slide_layouts):
        return prs.slide_layouts[spec.index]
    if spec.name:
        target = spec.name.strip().lower()
        for layout in prs.slide_layouts:
            if str(layout.name).strip().lower() == target:
                return layout
    return None


def _slide_role(spec: SlideSpec, index: int, total: int) -> str:
    if index == 1:
        return "cover"
    if index == total:
        return "closing"
    raw = (spec.layout or spec.template_slide_id or spec.slide_type or "content").strip().lower()
    aliases = {
        "toc": "toc",
        "contents": "toc",
        "agenda": "toc",
        "timeline": "timeline",
        "process": "timeline",
        "comparison": "comparison",
        "data": "data",
        "chart": "data",
        "section": "section",
    }
    if _first_image(spec) and raw == "content":
        return "content_image"
    return aliases.get(raw, raw or "content")


def _remove_template_seed_slides(prs: Presentation, count: int) -> None:
    for _ in range(min(count, len(prs.slides))):
        first_id = prs.slides._sldIdLst[0]
        r_id = first_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(first_id)


def _parse_layouts(raw: Any) -> dict[str, LayoutSpec]:
    if not isinstance(raw, dict):
        return {}
    output: dict[str, LayoutSpec] = {}
    for key, value in raw.items():
        normalized_key = str(key).strip().lower()
        if isinstance(value, str):
            output[normalized_key] = LayoutSpec(name=value)
        elif isinstance(value, int):
            output[normalized_key] = LayoutSpec(index=value)
        elif isinstance(value, dict):
            output[normalized_key] = LayoutSpec(
                name=str(value.get("layout_name") or value.get("name") or ""),
                index=_optional_int(value.get("index")),
                role=str(value.get("role") or ""),
            )
    return output


def _template_dirs() -> list[Path]:
    root = template_root()
    if not root.exists():
        return []
    return [path for path in sorted(root.iterdir()) if path.is_dir()]


def _template_config_path(template_dir: Path) -> Path | None:
    for name in ("template.json", "template.yaml", "template.yml"):
        path = template_dir / name
        if path.exists():
            return path
    return None


def _read_config(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8")
    if path.suffix.lower() == ".json":
        data = json.loads(text)
    else:
        data = yaml.safe_load(text)
    return deepcopy(data) if isinstance(data, dict) else {}


def _normalize_id(value: str) -> str:
    return value.strip().lower().replace("-", "_").replace(" ", "_")


def _list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _optional_int(value: Any) -> int | None:
    if value is None or value == "":
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None
